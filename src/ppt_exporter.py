"""
Generate a PowerPoint report from JIRA ticket data using a template.
"""
from __future__ import annotations

import shutil
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

from .client import JiraAttachment


TEMPLATE_PATH = Path(__file__).parent.parent / "PPTTemplate" / "Template.pptx"


# ── Helpers ──────────────────────────────────────────────────────────────────

def _fill_text_shape(shape: Any, text: str, color: str | None = None) -> None:
    """Fill a text shape, preserving the first run's formatting."""
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        run = para.runs[0]
        run.text = text
        if color:
            from pptx.oxml.ns import qn
            r_elem = run._r
            rPr = r_elem.find(qn("a:rPr"))
            if rPr is None:
                rPr = etree.SubElement(r_elem, qn("a:rPr"))
            for tag in (qn("a:uFillTx"), qn("a:uLnTx")):
                for elem in rPr.findall(tag):
                    rPr.remove(elem)
            for existing in rPr.findall(qn("a:solidFill")):
                rPr.remove(existing)
            solidFill = etree.SubElement(rPr, qn("a:solidFill"))
            srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgbClr.set("val", color)
    else:
        para.text = text


def _set_cell_text(tc: etree._Element, text: str, color: str | None = None) -> None:
    """
    Set the text of a table cell by finding the <a:r><a:t> chain inside
    <a:txBody><a:p> and replacing the <a:t> text. If no <a:r> exists,
    insert one.
    """
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        return

    def apply_color(r_elem: etree._Element, hex_color: str) -> None:
        """Add or replace <a:solidFill> inside the existing <a:rPr>, removing
        <a:uFillTx/> so the explicit color isn't overridden by the theme."""
        rPr = r_elem.find(qn("a:rPr"))
        if rPr is None:
            rPr = etree.SubElement(r_elem, qn("a:rPr"))
        # Remove theme-override flags that would supersede our solidFill
        for tag in (qn("a:uFillTx"), qn("a:uLnTx")):
            for elem in rPr.findall(tag):
                rPr.remove(elem)
        for existing in rPr.findall(qn("a:solidFill")):
            rPr.remove(existing)
        solidFill = etree.SubElement(rPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", hex_color)

    for para in txBody.findall(qn("a:p")):
        runs = para.findall(qn("a:r"))
        if runs:
            t_elem = runs[0].find(qn("a:t"))
            if t_elem is not None:
                t_elem.text = text
                if color:
                    apply_color(runs[0], color)
                return
        else:
            # No <a:r> in this paragraph — inject one before <a:endParaRPr>
            end_para = para.find(qn("a:endParaRPr"))
            new_r = etree.SubElement(para, qn("a:r"))
            new_t = etree.SubElement(new_r, qn("a:t"))
            new_t.text = text
            if color:
                apply_color(new_r, color)
            if end_para is not None:
                para.remove(end_para)
                para.append(end_para)
            return


def _fill_text_shape_mixed(
    shape: Any,
    segments: list[tuple[str, str | None]],
) -> None:
    """
    Fill a text shape with multiple segments, each optionally colored.
    Clones the first existing run's <a:rPr> (font properties) so that
    font size, bold, italic, etc. are preserved — only the color is
    overridden when specified.

    Args:
        shape: a PowerPoint shape with a text_frame.
        segments: list of (text, hex_color) tuples. hex_color is
                  a 6-char RGB hex string (e.g. "FF0000") or None
                  for the default color.
    """
    tf = shape.text_frame
    para = tf.paragraphs[0]
    p_elem = para._p

    # Clone the first existing <a:rPr> to preserve font size/weight/etc.
    existing_runs = p_elem.findall(qn("a:r"))
    orig_rPr = None
    if existing_runs:
        orig_r_elem = existing_runs[0]
        orig_rPr = etree.fromstring(etree.tostring(existing_runs[0].find(qn("a:rPr")))) if existing_runs[0].find(qn("a:rPr")) is not None else None

    # Remove all existing <a:r> elements from the paragraph
    for r_elem in list(p_elem.findall(qn("a:r"))):
        p_elem.remove(r_elem)

    # Insert runs before <a:endParaRPr> if present, otherwise at end
    end_para_rpr = p_elem.find(qn("a:endParaRPr"))

    for text, color in segments:
        r_elem = etree.SubElement(p_elem, qn("a:r"))

        if orig_rPr is not None:
            rPr_elem = etree.fromstring(etree.tostring(orig_rPr))
            # Strip theme-override flags so our solidFill isn't ignored
            for tag in (qn("a:uFillTx"), qn("a:uLnTx")):
                for elem in rPr_elem.findall(tag):
                    rPr_elem.remove(elem)
            r_elem.append(rPr_elem)

        if color:
            rPr_elem = r_elem.find(qn("a:rPr"))
            if rPr_elem is None:
                rPr_elem = etree.SubElement(r_elem, qn("a:rPr"))
            for existing in rPr_elem.findall(qn("a:solidFill")):
                rPr_elem.remove(existing)
            solidFill = etree.SubElement(rPr_elem, qn("a:solidFill"))
            srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgbClr.set("val", color)

        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text

        if end_para_rpr is not None:
            p_elem.remove(end_para_rpr)
            p_elem.append(end_para_rpr)


# ── Core class ────────────────────────────────────────────────────────────────

class PPTTemplateFiller:
    """
    Fills ``PPTTemplate/Template.pptx`` with JIRA ticket data.

    **Template contract (3 slides)**

    Slide 1 — Summary form
        * ``[Module | Ticket]`` table (headers only in the template).
          One row is added per exported ticket.
    Slide 2 — Diff per-ticket slide
        * ``[Module]`` text box
        * ``[key] [summary]`` text box
        * ``Diff`` label + picture placeholder
    Slide 3 — Self-Test per-ticket slide
        * ``[Module]`` text box
        * ``[key] [summary]`` text box
        * ``Self Test Report`` label + picture placeholder

    Output is written next to the JSON file: e.g. ``jira_export_pptx.pptx``.
    """

    def __init__(self, template_path: Path | str | None = None):
        if template_path is None:
            template_path = TEMPLATE_PATH
        self.template_path = Path(template_path)

    def fill(
        self,
        tickets: list[dict[str, Any]],
        output_path: Path | str | None = None,
        attachments_map: dict[str, list[JiraAttachment]] | None = None,
        session: Any = None,
    ) -> Path:
        """
        Generate the PPTX report.

        Args:
            tickets: list of ticket dicts from ``JiraExporter``.
                     Expected keys: ``key``, ``summary``, ``Module``,
                     ``Diff``, ``Self Test Report``.
            output_path: destination PPTX path. If ``None``,
                         derived from the template directory.
            attachments_map: optional ``{ticket_key: [JiraAttachment, ...]}``
                             for image replacement.
            session: authenticated ``requests.Session`` for downloading
                     JIRA attachment images.
        """
        if attachments_map is None:
            attachments_map = {}

        if not self.template_path.exists():
            raise FileNotFoundError(
                f"PPTX template not found: {self.template_path}. "
                "Please create a PPTTemplate/Template.pptx file or use --ppt-template to specify one."
            )

        self._session = session

        tmp_dir = Path(tempfile.mkdtemp())
        tmp_copy = tmp_dir / "Template.pptx"
        shutil.copy(self.template_path, tmp_copy)

        prs = Presentation(tmp_copy)
        if len(prs.slides) < 3:
            raise ValueError(
                f"Template must have at least 3 slides, found {len(prs.slides)}"
            )

        # ── Slide 1: summary table ──────────────────────────────────────────
        self._fill_summary_slide(prs.slides[0], tickets)

        # ── Fill the original template slides with the first ticket ──────────
        # Then duplicate them for remaining tickets
        if tickets:
            first = tickets[0]
            self._fill_ticket_slide(prs.slides[1], first, "Diff", attachments_map)
            self._fill_ticket_slide(prs.slides[2], first, "Self Test Report", attachments_map)

        # ── Per-ticket slides ───────────────────────────────────────────────
        # Slides 2-3 are filled with the first ticket above.
        # For tickets 2..N, duplicate those slides and fill them.
        for ticket in tickets[1:]:
            diff_slide = self._duplicate_slide(prs, 1)
            self._fill_ticket_slide(diff_slide, ticket, "Diff", attachments_map)

            str_slide = self._duplicate_slide(prs, 2)
            self._fill_ticket_slide(str_slide, ticket, "Self Test Report", attachments_map)

        # ── Save ───────────────────────────────────────────────────────────
        prs.save(str(tmp_copy))

        if output_path:
            final_path = Path(output_path)
        else:
            final_path = self.template_path.parent.parent / "jira_export_pptx.pptx"

        # On Windows, a locked file may need a short wait before overwriting
        import time
        for attempt in range(3):
            try:
                if final_path.exists():
                    final_path.unlink()
                shutil.copy(str(tmp_copy), str(final_path))
                break
            except OSError:
                if attempt < 2:
                    time.sleep(1)
                else:
                    raise
        shutil.rmtree(tmp_dir, ignore_errors=True)
        return final_path

    # ── Internal helpers ─────────────────────────────────────────────────────

    def _fill_summary_slide(
        self,
        slide: Any,
        tickets: list[dict[str, Any]],
    ) -> None:
        """Find the table on slide 1 and add one row per ticket."""
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            tbl = table._tbl

            # Clone row 1 (the first data row — empty in the template)
            template_row = tbl.tr_lst[1] if len(tbl.tr_lst) > 1 else tbl.tr_lst[0]

            for ticket in tickets:
                new_row = etree.fromstring(etree.tostring(template_row))
                cells = new_row.findall(qn("a:tc"))
                is_bug = ticket.get("issue_type", "").lower() == "bug"
                bug_color = "FF0000" if is_bug else None
                if len(cells) >= 2:
                    _set_cell_text(cells[0], ticket.get("Module", ""), color=bug_color)
                    _set_cell_text(cells[1], ticket.get("key", ""), color=bug_color)
                tbl.append(new_row)
            return

    def _duplicate_slide(self, prs: Presentation, src_index: int) -> Any:
        """
        Return a new slide that is a structural copy of slides[src_index].
        All shapes (text, picture, table, etc.) are cloned so the new slide
        looks identical — the caller can then mutate text and pictures.
        """
        src_slide = prs.slides[src_index]

        # Match the source layout by name
        src_name = src_slide.slide_layout.name
        layout = next(
            (l for l in prs.slide_layouts if l.name == src_name),
            prs.slide_layouts[0],
        )
        new_slide = prs.slides.add_slide(layout)

        # Remove the placeholder shapes that python-pptx inserts for us
        for sp in list(new_slide.shapes._spTree):
            tag = sp.tag.split("}")[-1] if "}" in sp.tag else sp.tag
            if tag in ("sp", "grpSp"):
                new_slide.shapes._spTree.remove(sp)

        # Clone every shape from the source slide
        for shape in src_slide.shapes:
            el = shape._element
            new_el = etree.fromstring(etree.tostring(el))
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        return new_slide

    def _fill_ticket_slide(
        self,
        slide: Any,
        ticket: dict[str, Any],
        role: str,  # "Diff" or "Self Test Report"
        attachments_map: dict[str, list[JiraAttachment]],
    ) -> None:
        """
        Fill placeholders on a duplicated per-ticket slide and replace images.
        """
        key = ticket.get("key", "")
        summary = ticket.get("summary", "")
        module = ticket.get("Module", "")
        is_bug = ticket.get("issue_type", "").lower() == "bug"
        bug_color = "FF0000" if is_bug else None
        key_summary = f"{key} {summary}"
        ticket_atts = attachments_map.get(key, [])

        for shape in slide.shapes:
            # Handle picture shapes separately
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self._replace_picture(slide, shape, ticket, role, ticket_atts)
                continue

            # Handle text shapes
            if not hasattr(shape, "text_frame"):
                continue

            raw = shape.text_frame.text.strip()

            if raw == "[Module]":
                _fill_text_shape(shape, module, color=bug_color)
            elif raw in ("[key] [summary]", "[key][summary]"):
                _fill_text_shape_mixed(
                    shape,
                    [(key, bug_color), (" " + summary, None)],
                )
            elif raw.startswith("[") and raw.endswith("]"):
                field_key = raw[1:-1]
                val = ticket.get(field_key, "")
                if isinstance(val, list):
                    val = ", ".join(str(v) for v in val)
                if val:
                    _fill_text_shape(shape, str(val))

    def _replace_picture(
        self,
        slide: Any,
        picture_shape: Any,
        ticket: dict[str, Any],
        role: str,
        ticket_atts: list[JiraAttachment],
    ) -> None:
        """
        Replace a picture placeholder identified by ``role`` with the image
        referenced in the ticket's ``Diff`` or ``Self Test Report`` field.
        """
        # Map the role label to the actual ticket field that holds the URL(s).
        # "Diff" images come from the "Solution" custom field; "Self Test Report"
        # maps to the field of the same name.
        field_name = "Solution" if role == "Diff" else role
        field_urls = ticket.get(field_name)
        if not field_urls:
            return

        url = str(field_urls[0]).strip() if isinstance(field_urls, list) else str(field_urls).strip()
        if not url.startswith("http"):
            return

        # Try to download via a matching attachment URL, then directly
        local_path = None
        for att in ticket_atts:
            if url in att.url or att.url in url:
                local_path = self._download_attachment(att)
                if local_path:
                    break

        if not local_path:
            local_path = self._download_image(url)

        if not local_path or not Path(local_path).exists():
            return

        from PIL import Image as PILImage
        img = PILImage.open(local_path)
        img_w, img_h = img.size
        img.close()

        max_w = picture_shape.width
        max_h = picture_shape.height
        scale = min(max_w / img_w, max_h / img_h)
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)

        left = picture_shape.left + (max_w - new_w) // 2
        top = picture_shape.top + (max_h - new_h) // 2

        # Remove old picture from XML tree
        picture_shape._element.getparent().remove(picture_shape._element)

        # Insert new picture with correct aspect ratio, centered in placeholder
        slide.shapes.add_picture(str(local_path), left, top, new_w, new_h)

    def _download_attachment(self, att: JiraAttachment) -> str | None:
        session = getattr(self, "_session", None) or self._make_session()
        try:
            resp = session.get(att.url, timeout=30)
            if resp.ok:
                ext = Path(att.filename).suffix or ".png"
                tmp = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
                tmp.write(resp.content)
                tmp.close()
                return tmp.name
        except Exception:
            pass
        return None

    def _download_image(self, url: str) -> str | None:
        session = getattr(self, "_session", None) or self._make_session()
        try:
            resp = session.get(url, timeout=30)
            if not resp.ok:
                return None
            ct = resp.headers.get("Content-Type", "")
            ext = ".png"
            if "jpeg" in ct or "jpg" in ct:
                ext = ".jpg"
            elif "gif" in ct:
                ext = ".gif"
            elif "webp" in ct:
                ext = ".webp"
            tmp = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
            tmp.write(resp.content)
            tmp.close()
            return tmp.name
        except Exception:
            pass
        return None

    def _make_session(self) -> Any:
        """Create a plain requests session (for backwards compatibility)."""
        import requests
        return requests
